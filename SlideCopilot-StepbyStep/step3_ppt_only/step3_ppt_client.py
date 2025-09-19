
"""
Step 3 — Generate a plain PowerPoint from JSON outline (no template)
- Reuses the RAG + LLM pipeline from Step 2
- Renders slides with python-pptx in a vanilla theme
"""

import os, io, json, re
from typing import List, Tuple

# ---------- RAG + LLM (copied/minified from Step 2) ----------
OPENAI_API_KEY = ""      # e.g., "sk-..."; leave empty to read from env
#OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-1.5-pro")

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

def _normalize_ws(t: str) -> str:
    return re.sub(r"\s+", " ", (t or "")).strip()

def build_corpus(texts: List[str]):
    clean = [_normalize_ws(x) for x in texts if x and x.strip()]
    if not clean:
        clean = [""]
    vec = TfidfVectorizer(max_features=8000)
    X = vec.fit_transform(clean)
    return vec, X, clean

def top_k(query: str, vec, X, corpus: List[str], k: int = 5) -> List[Tuple[int, float, str]]:
    qv = vec.transform([_normalize_ws(query)])
    sims = cosine_similarity(qv, X)[0]
    idxs = sims.argsort()[::-1][:k]
    return [(int(i), float(sims[i]), corpus[i]) for i in idxs]

def _extract_json_array(text: str) -> list:
    text = (text or "").strip()
    start = text.find("["); end = text.rfind("]")
    if start != -1 and end != -1 and end > start:
        return json.loads(text[start:end+1])
    if text.startswith("{"):
        obj = json.loads(text)
        if isinstance(obj, dict) and "slides" in obj:
            return obj["slides"]
    return json.loads(text)

def _fallback_outline(msg: str) -> list:
    return [{"layout":"title_content","title":"Fallback Outline","bullets":[msg], "notes":""}]

def _system_prompt(brand_rules: str) -> str:
    base = "You are a slide strategist. Output only a JSON array of slide specs; no prose."
    if brand_rules:
        base += " Brand rules: " + brand_rules
    return base

def _user_instructions(context: str, user_prompt: str) -> str:
    schema = """
Return a JSON list where each item is a slide spec:
- layout: one of ["title_content","two_content","chart"]
- title: string
- bullets: array of strings (for title_content)
- left_bullets/right_bullets: arrays (for two_content)
- chart: {categories:[], values:[], series_name:str} (for chart)
- notes: string
Keep bullets concise (<=12 words), <=6 bullets per slide.
Use the context faithfully. If context contradicts the prompt, prefer context.
"""
    return f"Context:\n{context}\n\nUser prompt:\n{user_prompt}\n\n{schema}"

def call_openai_outline(context: str, prompt: str, brand_rules: str = "") -> list:
    if not OPENAI_API_KEY:
        return _fallback_outline("Missing OPENAI_API_KEY")
    try:
        from openai import OpenAI
        client = OpenAI(api_key=OPENAI_API_KEY)
        sys = _system_prompt(brand_rules)
        user = _user_instructions(context, prompt)
        resp = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role":"system","content":sys},{"role":"user","content":user}],
            temperature=0.2,
        )
        text = (resp.choices[0].message.content or "").strip()
        return _extract_json_array(text)
    except Exception as e:
        return [{"layout":"title_content","title":"Model Error (OpenAI)","bullets":[type(e).__name__, str(e)[:180]],"notes":""}]

def call_gemini_outline(context: str, prompt: str, brand_rules: str = "") -> list:
    if not GEMINI_API_KEY:
        return _fallback_outline("Missing GEMINI_API_KEY")
    try:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel(GEMINI_MODEL)
        sys = _system_prompt(brand_rules)
        user = _user_instructions(context, prompt)
        resp = model.generate_content(f"{sys}\n\n{user}")
        text = (resp.text or "").strip()
        return _extract_json_array(text)
    except Exception as e:
        return [{"layout":"title_content","title":"Model Error (Gemini)","bullets":[type(e).__name__, str(e)[:180]],"notes":""}]

def generate_outline_with_rag(provider: str, prompt: str, brand_rules: str, texts: List[str], k: int = 5):
    vec, X, corpus = build_corpus(texts or [""])
    picks = top_k(prompt, vec, X, corpus, k=k)
    context = "\n\n---\n\n".join([p[2] for p in picks])
    if (provider or "").lower() == "openai":
        outline = call_openai_outline(context, prompt, brand_rules)
    elif (provider or "").lower() == "gemini":
        outline = call_gemini_outline(context, prompt, brand_rules)
    else:
        outline = call_openai_outline(context, prompt, brand_rules)
        if outline and isinstance(outline, list):
            pass
        else:
            outline = call_gemini_outline(context, prompt, brand_rules)
    return outline, picks

# ---------- PowerPoint rendering (no template) ----------
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

def _add_bullets(shape, items: List[str]):
    tf = shape.text_frame
    tf.clear()
    for i, text in enumerate(items or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0

def render_plain_ppt(outline: list, output_path: str) -> str:
    prs = Presentation()  # blank theme

    # grab common built-in layouts
    title_slide_layout = prs.slide_layouts[0]    # Title Slide
    title_and_content = prs.slide_layouts[1]     # Title and Content
    section_header    = prs.slide_layouts[2]     # Section Header
    two_content       = prs.slide_layouts[3]     # Two Content

    for spec in outline:
        layout = (spec.get("layout") or "title_content").lower()
        title = spec.get("title", "")

        if layout == "two_content":
            slide = prs.slides.add_slide(two_content)
            slide.shapes.title.text = title
            # placeholders: 1 and 2 are bodies
            try:
                left_ph = slide.placeholders[1]
                right_ph = slide.placeholders[2]
                _add_bullets(left_ph, spec.get("left_bullets"))
                _add_bullets(right_ph, spec.get("right_bullets"))
            except Exception:
                pass

        elif layout == "chart":
            slide = prs.slides.add_slide(title_and_content)
            slide.shapes.title.text = title
            # clear body area
            try:
                body = slide.placeholders[1]
                body.text_frame.clear()
            except Exception:
                pass
            ch = spec.get("chart", {})
            chart_data = ChartData()
            chart_data.categories = ch.get("categories", [])
            chart_data.add_series(ch.get("series_name", "Series"), ch.get("values", []))
            left, top, width, height = Inches(1), Inches(2.0), Inches(8), Inches(4)
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)

        else:
            slide = prs.slides.add_slide(title_and_content)
            slide.shapes.title.text = title
            try:
                body = slide.placeholders[1]
                _add_bullets(body, spec.get("bullets"))
            except Exception:
                pass

        # Speaker notes (if provided)
        if spec.get("notes"):
            slide.notes_slide.notes_text_frame.text = spec["notes"]

    prs.save(output_path)
    return output_path
