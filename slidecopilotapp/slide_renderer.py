
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

def _safe_get_title(slide):
    t = slide.shapes.title
    if t:
        return t
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph
    return None

def _find_body_placeholders(slide):
    bodies = []
    for ph in slide.placeholders:
        try:
            _ = ph.text_frame
            if ph.placeholder_format.idx != 0:
                bodies.append(ph)
        except Exception:
            pass
    return bodies

def _add_bullets(ph, bullets):
    tf = ph.text_frame
    tf.clear()
    for i, txt in enumerate(bullets or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = 0

def _set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text or ""

def render_from_outline(template_path, outline, output_path):
    prs = Presentation(template_path)

    def get_layout_by_name(name_lower):
        for l in prs.slide_layouts:
            if l.name.lower() == name_lower:
                return l
        return None
    def get_layout_with_min_placeholders(n=2):
        for l in prs.slide_layouts:
            if len(l.placeholders) >= n:
                return l
        return prs.slide_layouts[0]

    title_content = get_layout_by_name("title and content") or get_layout_with_min_placeholders(2)
    two_content   = get_layout_by_name("two content") or get_layout_with_min_placeholders(3)

    for spec in outline:
        layout = spec.get("layout", "title_content")
        if layout == "two_content":
            slide = prs.slides.add_slide(two_content)
            t = _safe_get_title(slide);  t.text = spec.get("title","") if t else None
            bodies = _find_body_placeholders(slide)
            if len(bodies) >= 2:
                _add_bullets(bodies[0], spec.get("left_bullets"))
                _add_bullets(bodies[1], spec.get("right_bullets"))
            _set_notes(slide, spec.get("notes"))
        elif layout == "chart":
            slide = prs.slides.add_slide(title_content)
            t = _safe_get_title(slide);  t.text = spec.get("title","") if t else None
            bodies = _find_body_placeholders(slide)
            if bodies:
                bodies[0].text_frame.clear()
            ch = spec.get("chart", {})
            chart_data = ChartData()
            chart_data.categories = ch.get("categories", [])
            chart_data.add_series(ch.get("series_name","Series"), ch.get("values", []))
            left, top, width, height = Inches(1), Inches(2.0), Inches(8), Inches(4)
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
            _set_notes(slide, spec.get("notes"))
        else:
            slide = prs.slides.add_slide(title_content)
            t = _safe_get_title(slide);  t.text = spec.get("title","") if t else None
            bodies = _find_body_placeholders(slide)
            if bodies:
                _add_bullets(bodies[0], spec.get("bullets"))
            _set_notes(slide, spec.get("notes"))

    prs.save(output_path)
    return output_path
