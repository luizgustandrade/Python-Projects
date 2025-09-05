
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

def get_layout_by_name(prs, target_name_lower):
    for layout in prs.slide_layouts:
        if layout.name.lower() == target_name_lower:
            return layout
    return None

def get_layout_with_min_placeholders(prs, min_count=2):
    for layout in prs.slide_layouts:
        if len(layout.placeholders) >= min_count:
            return layout
    return prs.slide_layouts[0]

def safe_get_title(slide):
    t = slide.shapes.title
    if t:
        return t
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph
    return None

def find_body_placeholders(slide):
    bodies = []
    for ph in slide.placeholders:
        try:
            _ = ph.text_frame
            if ph.placeholder_format.idx != 0:
                bodies.append(ph)
        except Exception:
            pass
    return bodies

def add_bullets_to_placeholder(ph, bullets):
    tf = ph.text_frame
    tf.clear()
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = 0

def notes_set(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def render_title_content(prs, layout, title, bullets=None, notes=None):
    slide = prs.slides.add_slide(layout)
    t = safe_get_title(slide)
    if t:
        t.text = title or ""
    bodies = find_body_placeholders(slide)
    if bullets and bodies:
        add_bullets_to_placeholder(bodies[0], bullets)
    if notes:
        notes_set(slide, notes)
    return slide

def render_two_content(prs, layout, title, left_bullets=None, right_bullets=None, notes=None):
    slide = prs.slides.add_slide(layout)
    t = safe_get_title(slide)
    if t:
        t.text = title or ""
    bodies = find_body_placeholders(slide)
    if len(bodies) >= 2:
        if left_bullets:
            add_bullets_to_placeholder(bodies[0], left_bullets)
        if right_bullets:
            add_bullets_to_placeholder(bodies[1], right_bullets)
    if notes:
        notes_set(slide, notes)
    return slide

def render_chart_slide(prs, layout, title, categories, values, series_name="Series", notes=None):
    slide = prs.slides.add_slide(layout)
    t = safe_get_title(slide)
    if t:
        t.text = title or ""
    bodies = find_body_placeholders(slide)
    if bodies:
        bodies[0].text_frame.clear()
    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    left, top, width, height = Inches(1), Inches(2.0), Inches(8), Inches(4)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
    if notes:
        notes_set(slide, notes)
    return slide

def build_deck(template_path, output_path):
    prs = Presentation(template_path)

    title_content_layout = get_layout_by_name(prs, "title and content") or get_layout_with_min_placeholders(prs, 2)
    title_only_layout    = get_layout_by_name(prs, "title only") or get_layout_with_min_placeholders(prs, 1)

    two_content_layout = get_layout_by_name(prs, "two content")
    if two_content_layout is None:
        candidate = None
        for layout in prs.slide_layouts:
            if len(layout.placeholders) >= 3:
                candidate = layout
                break
        two_content_layout = candidate if candidate else title_content_layout

    # Example outline (swap for LLM output in your app)
    render_title_content(
        prs, title_content_layout,
        title="AI Sales Masterclass – Highlights",
        bullets=["What’s new in GenAI Services", "Customer wins & repeatable plays", "Next steps for the field"],
        notes="30–40s overview. Mention LATAM best practices.",
    )

    render_two_content(
        prs, two_content_layout,
        title="Pipeline Levers",
        left_bullets=["Workstation attach", "GPU deals → Services attach", "Partner enablement"],
        right_bullets=["Workshops", "Pilot scope kits", "Customer references"],
        notes="Call out 20% YoY pipeline lift target.",
    )

    render_chart_slide(
        prs, title_content_layout,
        title="Q3 Focus Metrics",
        categories=["Jul", "Aug", "Sep"],
        values=[5, 11, 16],
        series_name="Workshops",
        notes="Tie to regional goals.",
    )

    prs.save(output_path)

if __name__ == "__main__":
    TEMPLATE = "Template-PPT.pptx"  # put in same folder
    OUT = "generated_deck_from_template.pptx"
    build_deck(TEMPLATE, OUT)
    print("Saved:", OUT)
